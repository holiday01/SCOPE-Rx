"""Build Phase-1 PPT deck from figures + results."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import pandas as pd, json

ROOT = Path('/home/holiday01/drug_sc')
FIG  = ROOT/'results/figures'
OUT  = ROOT/'results/SCOPE-Rx_Phase1_report.pptx'

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

ACCENT = RGBColor(0x26, 0x46, 0x53)
SUB    = RGBColor(0x55, 0x55, 0x55)
GOOD   = RGBColor(0x2A, 0x9D, 0x8F)
BAD    = RGBColor(0xE7, 0x6F, 0x51)

def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_title(slide, text, size=28):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = True; p.font.color.rgb = ACCENT

def add_subtitle(slide, text, top_in=1.1, size=14):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(top_in), Inches(12.5), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.color.rgb = SUB; p.font.italic = True

def add_bullets(slide, bullets, left=0.4, top=1.6, width=6.5, height=5.5, size=14):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        if isinstance(text, tuple):
            bold, rest = text
            r1 = p.add_run(); r1.text = bold + ' — '; r1.font.bold = True; r1.font.size = Pt(size)
            r2 = p.add_run(); r2.text = rest; r2.font.size = Pt(size)
        else:
            p.text = '• ' + text
            p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(0x22,0x22,0x22)
        p.space_after = Pt(6)

def add_image(slide, path, left_in, top_in, width_in=None, height_in=None):
    kw = {}
    if width_in: kw['width']=Inches(width_in)
    if height_in: kw['height']=Inches(height_in)
    slide.shapes.add_picture(str(path), Inches(left_in), Inches(top_in), **kw)

# =================== Slide 1: Title ===================
s = blank_slide()
add_title(s, 'SCOPE-Rx — Phase-1 Report', size=40)
tb = s.shapes.add_textbox(Inches(0.4), Inches(1.1), Inches(12.5), Inches(1.0))
p = tb.text_frame.paragraphs[0]
p.text = 'Multi-modal single-cell × drug-structure × TCGA survival pipeline for HCC drug repositioning'
p.font.size = Pt(20); p.font.color.rgb = SUB; p.font.italic = True

# stats block
stats = [
    ('Pilot cancer', 'LIHC / Hepatocellular carcinoma'),
    ('Bulk data', 'TCGA-LIHC — 423 patients × 20,530 genes + OS clinical'),
    ('Single-cell', 'Integrated HCC atlas — 172,538 cells × 29,834 genes, 6 GEO series'),
    ('Drug space', 'DepMap + GDSC1+2 + PRISM + CTRPv2 — 1,806 drugs × 25 HCC lines, 1.13 M (line×drug) pairs'),
    ('DL stack', 'PyTorch 2.10 + cu128 · Geneformer V2-104M (768-d) · MPNN (rdkit+PyG) · Attention deconv · Cox'),
    ('Hardware', 'RTX 5070 (Blackwell sm_120, 16 GB) · ≈ 8 min total training time'),
]
tb = s.shapes.add_textbox(Inches(0.4), Inches(2.3), Inches(12.5), Inches(4.5))
tf = tb.text_frame
for i,(k,v) in enumerate(stats):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    r1 = p.add_run(); r1.text = k + ':  '; r1.font.bold = True; r1.font.size = Pt(15); r1.font.color.rgb = ACCENT
    r2 = p.add_run(); r2.text = v; r2.font.size = Pt(15); r2.font.color.rgb = RGBColor(0x33,0x33,0x33)
    p.space_after = Pt(8)

tb = s.shapes.add_textbox(Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.5))
p = tb.text_frame.paragraphs[0]
p.text = 'Generated 2026-04-24'
p.font.size = Pt(11); p.font.color.rgb = SUB; p.font.italic = True

# =================== Slide 2: Pipeline overview ===================
s = blank_slide()
add_title(s, 'Pipeline overview')
add_subtitle(s, 'Six tasks from raw data to wet-lab-ready drug list, all patches applied')
add_image(s, FIG/'fig01_pipeline.png', 0.4, 1.6, width_in=12.5)
add_bullets(s, [
    ('Inputs', 'TCGA bulk + scRNA atlas + pan-cancer cell-line AUC'),
    ('Cell state', 'Geneformer (768-d) → Leiden 57 prototypes with trust scores'),
    ('Scoring', 'Survival-anchored AUC ×  oncology filter  ×  pathway-Cox prior')
], top=6.1, size=12)

# =================== Slide 3: Baseline audit ===================
s = blank_slide()
add_title(s, 'Baseline audit — what works, what does not')
add_subtitle(s, 'Three independent designs; only scDEAL + T3c attention deconv survived')
add_image(s, FIG/'fig02_baselines_fix.png', 0.4, 1.6, width_in=12.5)
add_bullets(s, [
    ('scDEAL (gene-MLP + DANN)', 'per-line Spearman 0.61 — strongest single-drug head'),
    ('scPDS (pathway transformer)', '0.11 — pathway mean drops signal, transformer under-fits'),
    ('Scaden-CA', 'composition collapses to SNU886 (91 %) — 25 cell lines too narrow'),
    ('Cox HR bug', 'HRs were ~1e28 before compositional z-score; now 0.8–1.3 per SD')
], top=5.8, size=12)

# =================== Slide 4: Geneformer embedding ===================
s = blank_slide()
add_title(s, 'Cell-state foundation model — Geneformer V2-104M')
add_subtitle(s, 'Zero-shot 768-d embedding of 20k HCC cells, 10-NN label purity 0.80')
add_image(s, FIG/'fig09_geneformer_umap.png', 0.4, 1.6, height_in=5.3)
add_bullets(s, [
    ('Inference cost', '176 s on RTX 5070 fp16 for 20k cells (seq 2048)'),
    ('Usable genes', '17,642 / 29,834 (59 %) after mapping to gc104M token dictionary'),
    ('Downstream', 'prototype centroids (T3c) carry Geneformer embedding for neighbour queries'),
    ('Top cell-type purity', 'Neutrophil 0.92, Endothelial 0.90, Memory B 0.87, M2 macro 0.87'),
    ('Bug fix applied', 'use a.layers["counts"] not a.X (h5ad was pre-log-normalised)')
], left=9.0, top=1.6, width=4.2, size=11)

# =================== Slide 5: TCGA composition ===================
s = blank_slide()
add_title(s, 'TCGA deconvolution — attention beats linear mixing')
add_subtitle(s, 'Each patient decomposed across 55 HCC cell-state prototypes (no collapse)')
add_image(s, FIG/'fig03_composition.png', 0.4, 1.6, width_in=12.5)
add_bullets(s, [
    ('Entropy', '88.9 % of theoretical maximum — Scaden-CA was ~30 %'),
    ('Top-1 > 50 %', '0 patients collapse — Scaden-CA had ~60 % collapse'),
    ('Training', '60k Dirichlet-mixed pseudobulks from prototypes + 10 % noise, val MAE 0.009')
], top=5.8, size=12)

# =================== Slide 6: Prognostic prototypes ===================
s = blank_slide()
add_title(s, 'Prognosis-driving cell-state prototypes (Cox per SD)')
add_subtitle(s, 'TME signatures match HCC immunology literature')
add_image(s, FIG/'fig04_cox_prototypes.png', 0.4, 1.6, width_in=7.5)
add_bullets(s, [
    ('Bad prognosis (HR > 1)', 'CD14 macrophage (1.26), CD4 (1.24), Endothelial (1.23), M2 macrophage (1.19), CD16 macrophage (1.19)'),
    ('Protective (HR < 1)', 'Epithelial "normal" (0.82), CD4 effector (0.80), NK / NKT (0.81)'),
    ('Clinical match', 'TAM-high = poor OS (Zhang 2019, Sun 2021); NK-high = better OS'),
    ('Asymmetric weighting', 'bad × 1.0, good × 0.3 — avoids over-penalising killing protective subpops')
], left=8.1, top=1.6, width=5.0, size=11)

# =================== Slide 7: Pathway-Cox prior (Patch 1) ===================
s = blank_slide()
add_title(s, 'Patch 1 — drug-target × TCGA pathway-Cox prior')
add_subtitle(s, 'MSigDB Hallmark + KEGG + Reactome → 1,790 pathways × TCGA OS Cox')
add_image(s, FIG/'fig07_pathway_volcano.png', 0.4, 1.6, width_in=7.0)
add_bullets(s, [
    ('Hazard pathways', 'Hyaluronan metabolism (CD44), Chk1/Chk2, Nucleotide/PPP biosynthesis, ERBB4 signalling, ROS'),
    ('Protective', 'Urea cycle, Bile-acid synthesis, Lipoprotein remodelling — normal hepatocyte metabolism'),
    ('Drug–pathway links', '64,105 non-zero coverage cells; 1,111/1,489 drugs hit ≥1 hazard pathway'),
    ('Formula', 'S_prior(d) = Σ_pw target_cov(d,pw) × (-log10 p_pw) × sign(logHR_pw)')
], left=7.5, top=1.6, width=5.6, size=11)

# =================== Slide 8: Oncology filter (Patch 2) ===================
s = blank_slide()
add_title(s, 'Patch 2 — oncology / phase / MOA filter')
add_subtitle(s, 'Systematic rule-based re-weighting using PRISM metadata')
rows = [
    ('+3', 'disease.area contains oncology / cancer / tumour / leukaemia / lymphoma …'),
    ('+2', 'indication mentions cancer-type keyword'),
    ('+2', 'any target is in COSMIC tier-1 cancer-driver list'),
    ('+1', 'phase == Launched'),
    ('-3', 'MOA matches non-oncology class (antimicrobial / antiseptic / GABA / anticonvulsant / …)'),
    ('-2', 'phase == Preclinical (no human safety data)'),
]
tb = s.shapes.add_textbox(Inches(0.4), Inches(1.6), Inches(12.5), Inches(3.5))
tf = tb.text_frame
for i,(w,d) in enumerate(rows):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    r1 = p.add_run(); r1.text = f'  {w}  '; r1.font.bold = True; r1.font.size = Pt(18)
    r1.font.color.rgb = GOOD if w.startswith('+') else BAD
    r2 = p.add_run(); r2.text = d; r2.font.size = Pt(14)
    p.space_after = Pt(6)

add_bullets(s, [
    ('Result', '1,602 / 1,806 drugs have PRISM metadata; 185 score ≥ 3, 39 blacklisted'),
    ('Strict filter', '777 drugs (oncology-tagged + Phase 1+)'),
    ('Soft score', 'z(kill) + 0.5·z(oncology) + 0.7·z(prior)'),
    ('Drugs removed', 'Cetrimonium, Alexidine, Tiagabine, Brivaracetam, Istradefylline, …')
], top=5.1, size=12)

# =================== Slide 9: Trust fix (Patch 3) ===================
s = blank_slide()
add_title(s, 'Patch 3 — honest cell-type-matched trust references')
add_subtitle(s, '≈350 DepMap lymphoid/myeloid lines matched by dominant cell type')
add_image(s, FIG/'fig08_trust_fix.png', 0.4, 1.6, width_in=7.0)
add_bullets(s, [
    ('Previous', 'trust = max Pearson r vs ANY DepMap line → inflated by housekeeping'),
    ('Now', 'match M2 macro → Myeloid/AML (80 lines); T cells → T-ALL (26); NK → NK (4); B → B-cell (121)'),
    ('Finding', 'TME trust drops from 0.4 → 0.25 — honest, not worse'),
    ('Interpretation', 'DepMap has no true tissue macrophage → AUC extrapolation there is inherently limited'),
    ('Endothelial', 'fallback full-panel — requires external HUVEC for Phase-2')
], left=7.5, top=1.6, width=5.6, size=11)

# =================== Slide 10: Clinical sanity ===================
s = blank_slide()
add_title(s, 'Clinical HCC drug ranking — progression across patches')
add_subtitle(s, 'Known HCC drugs rising toward top after each correction')
add_image(s, FIG/'fig06_clinical_progression.png', 0.4, 1.5, height_in=5.7)
add_bullets(s, [
    ('Top 3 % (final)', 'Lapatinib #1, Afatinib #11, Erlotinib #26'),
    ('Top 5–10 %', 'Paclitaxel #40, Doxorubicin #198, Vincristine #142'),
    ('Top 15 %', 'Oxaliplatin, Lenvatinib, Regorafenib, 5-FU'),
    ('Residual miss', 'Sorafenib #467 — benefit is angiogenic/microenvironment, invisible at TCGA bulk transcription')
], left=8.2, top=1.6, width=4.9, size=11)

# =================== Slide 11: Top-20 decomposition ===================
s = blank_slide()
add_title(s, 'Top-20 drugs — decomposition of final score')
add_subtitle(s, '1.0·z(kill) + 0.5·z(oncology) + 0.7·z(pathway prior)')
add_image(s, FIG/'fig05_top20_decomp.png', 0.4, 1.6, height_in=5.5)
add_bullets(s, [
    ('Convergent themes', 'EGFR/HER (×6), Proteasome (×3), HSP90 (×2), AKT/PI3K (×4), CDK/MEK (×3)'),
    ('Highest kill-score', 'Lapatinib, AZD8931, PHA-767491, Afatinib'),
    ('Highest prior-score', 'Bortezomib, Carfilzomib, CID-5458317, HSP90 inhibitors'),
    ('Likely noise', 'CID-5458317 (preclinical only) — high prior but no clinical data')
], left=9.2, top=1.6, width=3.9, size=11)

# =================== Slide 12: Mechanism pie ===================
s = blank_slide()
add_title(s, 'Mechanistic convergence in Top-25')
add_subtitle(s, 'Pipeline picks oncology-validated pathways, not random cytotoxins')
add_image(s, FIG/'fig10_mechanism_pie.png', 2.0, 1.6, height_in=5.6)
add_bullets(s, [
    ('EGFR / HER family', 'Lapatinib, Afatinib, Gefitinib, Osimertinib, Erlotinib, AZD8931'),
    ('Proteasome', 'Bortezomib, Carfilzomib'),
    ('HSP90', 'Alvespimycin, Tanespimycin'),
    ('AKT / PI3K', 'MK-2206, Uprosertib, Copanlisib, Taselisib, Hexamethylenebisacetamide'),
    ('CDK / MEK', 'BMS-265246, Trametinib, PHA-767491'),
    ('Consistency', 'All five classes are active or trialled in HCC literature')
], left=7.6, top=1.6, width=5.5, size=12)

# =================== Slide 13: Wet-lab brief (Top-5) ===================
s = blank_slide()
add_title(s, 'Wet-lab brief — Top-5 candidates (sample)')
add_subtitle(s, 'Each drug carries target subpops, FACS markers, cell-line match, SMILES')

brief = [
    {'rank':1,'drug':'Lapatinib','moa':'EGFR / HER2 TKI','phase':'Launched',
     'sub':'#43 NK/NKT · #46 CD16 macro · #0 M2 macro',
     'markers':'NKG7 / GNLY / APOA1 / HLA-DRA / FTL',
     'line':'HUH6 (r=0.24) / SNU761 (r=0.18)'},
    {'rank':2,'drug':'Bortezomib','moa':'Proteasome / NF-κB','phase':'Launched',
     'sub':'#6 CD8 T · #18 CD4 · #21 M2 macro',
     'markers':'C1QA / HSPH1 / CTSB / HSP90AB1',
     'line':'SNU398 (r=0.17)'},
    {'rank':3,'drug':'Alvespimycin','moa':'HSP90','phase':'Phase 2',
     'sub':'#46 CD16 macro · #48 M2 macro · #43 NK/NKT',
     'markers':'C1QA / HLA-DPA1 / GNLY / HSPA1A',
     'line':'SNU761 (r=0.18)'},
    {'rank':4,'drug':'Afatinib','moa':'pan-HER TKI','phase':'Launched',
     'sub':'#22 Epithelial_tumor · #26 Epithelial_tumor',
     'markers':'APOA1 / ALB / AFP · AFP+ FACS gate',
     'line':'HUH6 (r=0.24)'},
    {'rank':5,'drug':'MK-2206','moa':'allosteric AKT','phase':'Phase 2',
     'sub':'#22 Epithelial_tumor · #26 Epithelial_tumor',
     'markers':'APOA1 / APOA2 / ALB',
     'line':'SNU398 (r=0.17)'},
]
tb = s.shapes.add_textbox(Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.5))
tf = tb.text_frame
for i,r in enumerate(brief):
    if i>0: tf.add_paragraph().space_after = Pt(4)
    p = tf.paragraphs[-1] if i==0 else tf.add_paragraph()
    r1 = p.add_run(); r1.text = f"  #{r['rank']}  {r['drug']}  "; r1.font.bold=True; r1.font.size=Pt(16); r1.font.color.rgb=ACCENT
    r2 = p.add_run(); r2.text = f"—  {r['moa']}   ({r['phase']})"; r2.font.size=Pt(14); r2.font.color.rgb=SUB
    p2 = tf.add_paragraph()
    r3 = p2.add_run(); r3.text = f"     Target subpops: {r['sub']}"; r3.font.size=Pt(11)
    p3 = tf.add_paragraph()
    r4 = p3.add_run(); r4.text = f"     FACS markers: {r['markers']}"; r4.font.size=Pt(11)
    p4 = tf.add_paragraph()
    r5 = p4.add_run(); r5.text = f"     Suggested line: {r['line']}"; r5.font.size=Pt(11)
    p4.space_after = Pt(4)

# =================== Slide 14: Known limitations ===================
s = blank_slide()
add_title(s, 'Known limitations — what this pipeline cannot see')
lims = [
    ('Sorafenib middle rank (467/1806)',
     'DepMap HCC lines are relatively resistant to sorafenib; clinical benefit is microenvironmental (VEGFR/RAF antiangiogenic). Bulk transcription misses this.'),
    ('TME trust stuck at 0.25',
     'No DepMap macrophage/T tissue reference. AML blasts ≠ Kupffer/TAM biology. Real validation requires PDO co-culture.'),
    ('Epithelial_normal vs Epithelial_tumor uses sample_type only',
     'No per-cell CNV inference yet — malignant/normal label can leak. Phase-2 should add infercnv.'),
    ('High-prior-preclinical drugs (e.g. CID-5458317)',
     'High pathway prior but zero clinical data. Flag, do not recommend for PDO until literature check.'),
    ('Unseen-drug MPNN (T3a) not integrated',
     '+0.30 Spearman on truly unseen compounds, but uses single joint head; Phase-2 should retrain as direction-vector generator.'),
    ('Endothelial reference missing',
     'DepMap has no HUVEC-class line. Phase-2: add ENCODE/Tabula-Sapiens endothelial bulk.'),
]
tb = s.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(12.5), Inches(6))
tf = tb.text_frame
for i,(a,b) in enumerate(lims):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    r1 = p.add_run(); r1.text = '• ' + a; r1.font.bold=True; r1.font.size=Pt(14); r1.font.color.rgb=BAD
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = '    ' + b; r2.font.size=Pt(12); r2.font.color.rgb=SUB
    p2.space_after = Pt(6)

# =================== Slide 15: Phase-2 roadmap ===================
s = blank_slide()
add_title(s, 'Phase-2 roadmap')
roadmap = [
    ('Short (1-2 wk, in silico)', [
        'Add CRISPR dependency prior (Achilles) — captures on-target essentiality beyond AUC',
        'Integrate MPNN T3a as direction-vector generator; extend to DrugBank/ChEMBL novel compounds',
        'Add infercnv-based malignant/normal hepatocyte split',
        'Add ENCODE HUVEC + DICE immune bulk as true TME references'
    ]),
    ('Medium (1 mo)', [
        'Expand to pan-cancer — HNSC, PAAD, BRCA as generalisation tests',
        'Replace prototype attention with full 172k scRNA softmax + sparsity prior',
        'Spatial transcriptomics module on Visium HCC cohort for niche-level predictions'
    ]),
    ('Wet-lab validation target', [
        'Top-5 candidates: Lapatinib, Bortezomib, Alvespimycin, Afatinib, MK-2206',
        'Assays: HCC PDO viability + FACS-sorted subpop killing + Perturb-seq signature matching',
        'Expected go/no-go within 6 weeks of organoid delivery'
    ]),
]
tb = s.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(12.5), Inches(6))
tf = tb.text_frame
for i,(h, bs) in enumerate(roadmap):
    if i>0: tf.add_paragraph().space_after = Pt(6)
    p = tf.paragraphs[-1] if i==0 else tf.add_paragraph()
    r = p.add_run(); r.text = h; r.font.bold=True; r.font.size=Pt(16); r.font.color.rgb=ACCENT
    for it in bs:
        p2 = tf.add_paragraph()
        p2.text = '    • ' + it
        p2.font.size = Pt(13); p2.space_after = Pt(2)

# =================== Slide 16: Summary ===================
s = blank_slide()
add_title(s, 'Summary')
add_subtitle(s, 'Every number the pipeline reports can be traced to a parquet file in results/')
summary_stats = [
    ('423', 'TCGA-LIHC patients with cell-state composition + top-5 drug ranking'),
    ('1,806', 'drugs ranked; 1,603 with SMILES, 1,488 with MOA/target/phase'),
    ('57', 'cell-state prototypes; 8 significant prognostic (Cox p<0.05)'),
    ('1,790', 'pathways tested in TCGA-LIHC Cox; 536 significant'),
    ('20,000', 'HCC cells embedded by Geneformer V2-104M; 10-NN purity 0.80'),
    ('5', 'mechanistic classes converge in Top-25 (EGFR, proteasome, HSP90, AKT/PI3K, CDK/MEK)'),
    ('3', 'patches fixed bug/prior/trust — each script re-run post-fix'),
    ('11 / 13', 'clinical HCC drugs now in top 25 %'),
]
tb = s.shapes.add_textbox(Inches(0.4), Inches(1.6), Inches(12.5), Inches(5))
tf = tb.text_frame
for i,(v,d) in enumerate(summary_stats):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    r1 = p.add_run(); r1.text = f'  {v:>7s}  '; r1.font.bold=True; r1.font.size=Pt(22); r1.font.color.rgb=ACCENT
    r2 = p.add_run(); r2.text = d; r2.font.size=Pt(14)
    p.space_after = Pt(8)

tb = s.shapes.add_textbox(Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.6))
p = tb.text_frame.paragraphs[0]
p.text = 'Deliverables:  results/t4/hcc_top20_wetlab_brief.md · results/t4/per_patient_top5.parquet · results/phase1_final_report.md'
p.font.size = Pt(11); p.font.italic=True; p.font.color.rgb = SUB

prs.save(OUT)
print(f'Saved: {OUT}')
print(f'Size: {OUT.stat().st_size/1024/1024:.2f} MB')
print(f'Slides: {len(prs.slides)}')
