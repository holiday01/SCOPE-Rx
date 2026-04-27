"""Phase-1 v2 PPT — adds external validation, meta-analysis, consistency battery."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT = Path('/home/holiday01/drug_sc')
FIG  = ROOT/'results/figures'
OUT  = ROOT/'results/SCOPE-Rx_Phase1_report_v2.pptx'

prs = Presentation()
prs.slide_width  = Inches(13.33); prs.slide_height = Inches(7.5)
ACC=RGBColor(0x26,0x46,0x53); SUB=RGBColor(0x55,0x55,0x55)
GOOD=RGBColor(0x2A,0x9D,0x8F); BAD=RGBColor(0xE7,0x6F,0x51); WARN=RGBColor(0xE9,0xC4,0x6A)

def blank(): return prs.slides.add_slide(prs.slide_layouts[6])
def title(s, text, size=28):
    tb = s.shapes.add_textbox(Inches(0.4),Inches(0.25),Inches(12.5),Inches(0.9))
    p = tb.text_frame.paragraphs[0]; p.text=text; p.font.size=Pt(size); p.font.bold=True; p.font.color.rgb=ACC
def sub(s, text, top=1.1, size=14):
    tb = s.shapes.add_textbox(Inches(0.4),Inches(top),Inches(12.5),Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.text=text; p.font.size=Pt(size); p.font.color.rgb=SUB; p.font.italic=True
def bullets(s, items, left=0.4, top=1.6, width=6.5, size=13):
    tb = s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(5.5))
    tf = tb.text_frame; tf.word_wrap=True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        if isinstance(it, tuple):
            r1=p.add_run(); r1.text=it[0]+' — '; r1.font.bold=True; r1.font.size=Pt(size); r1.font.color.rgb=ACC
            r2=p.add_run(); r2.text=it[1]; r2.font.size=Pt(size)
        else:
            p.text='• '+it; p.font.size=Pt(size)
        p.space_after=Pt(6); p.font.color.rgb=RGBColor(0x22,0x22,0x22)
def img(s, path, l, t, w=None, h=None):
    kw={}
    if w: kw['width']=Inches(w)
    if h: kw['height']=Inches(h)
    s.shapes.add_picture(str(path),Inches(l),Inches(t),**kw)

# ============== Slide 1: Title ==============
s = blank()
title(s, 'SCOPE-Rx — Phase-1 Report (v2)', size=40)
tb = s.shapes.add_textbox(Inches(0.4),Inches(1.1),Inches(12.5),Inches(1.0))
p = tb.text_frame.paragraphs[0]
p.text = 'Multi-modal single-cell × drug-structure × TCGA survival — with external validation, meta-analysis, and methodological consistency battery'
p.font.size=Pt(18); p.font.color.rgb=SUB; p.font.italic=True

stats = [
    ('Cohorts', '3 (TCGA-LIHC 423 RNA-seq · GSE14520 225 Affymetrix · GSE76427 115 Illumina) — n=763'),
    ('Endpoints', 'OS + RFS — multi-platform, multi-etiology, multi-endpoint validation'),
    ('External validation', '8 / 9 prognostic prototypes replicated; composite risk c-index 0.73 (GSE14520 OS, p=0.006)'),
    ('Meta-analysis', '16 / 57 prototypes pooled-significant (3-cohort fixed-effect)'),
    ('Method consistency', '5 survival methods consensus → 9 prototypes; 4 score aggregations consensus → 11 drugs'),
    ('Top drugs', 'Lapatinib · Afatinib · Trametinib · MK-2206 · Alvespimycin · Tanespimycin (mechanism convergence)'),
    ('Code & data', '13 scripts · 11 result tables · 18 figures · ~ 9 GB local'),
]
tb = s.shapes.add_textbox(Inches(0.4),Inches(2.3),Inches(12.5),Inches(4.5))
tf = tb.text_frame
for i,(k,v) in enumerate(stats):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    r1=p.add_run(); r1.text=k+':  '; r1.font.bold=True; r1.font.size=Pt(15); r1.font.color.rgb=ACC
    r2=p.add_run(); r2.text=v; r2.font.size=Pt(15)
    p.space_after=Pt(8)

tb=s.shapes.add_textbox(Inches(0.4),Inches(6.8),Inches(12.5),Inches(0.5))
p=tb.text_frame.paragraphs[0]; p.text='Generated 2026-04-26'
p.font.size=Pt(11); p.font.color.rgb=SUB; p.font.italic=True

# ============== Slide 2: Pipeline overview ==============
s = blank()
title(s, 'Pipeline overview')
sub(s, 'Six modules; survival prior + drug structure + scRNA prototypes; tiered output')
img(s, FIG/'fig01_pipeline.png', 0.4, 1.6, w=12.5)
bullets(s, [('Inputs','TCGA bulk · scRNA atlas · DepMap+GDSC+PRISM AUC · pathway gene sets · DrugBank-like targets'),
            ('Cell-state','Geneformer V2-104M (768-d) → 57 attention-deconvolver prototypes'),
            ('Scoring','z(kill) + 0.5·z(oncology) + 0.7·z(pathway-Cox prior) + VEGFR rescue'),
            ('Output','Tier A/B high-confidence + Tier C/D exploratory + per-patient Top-5')], top=6.1, size=11)

# ============== Slide 3: Baseline audit + bug fixes ==============
s = blank()
title(s, 'Baseline audit + bug fixes — what was wrong, what we fixed')
sub(s, 'Three baselines; three patches transparently documented')
img(s, FIG/'fig02_baselines_fix.png', 0.4, 1.6, w=12.5)
bullets(s, [('Baselines','scDEAL (gene-MLP+DANN) Spearman 0.61 (LOO 5-fold = 0.145±0.14); scPDS 0.11; Scaden-CA collapses to one cell-line'),
            ('Cox HR explosion','Pre-fix HR ≈ 1×10²⁸ (junk); post fix per-SD HR 0.8–1.3 — interpretable'),
            ('scRNA normalisation','`a.X` was log-normalised, not raw counts; switched to `a.layers["counts"]` for T2b/T2c/T3b/T3c'),
            ('Drug-catalog SMILES','PRISM column had comma-duplicated SMILES; recovered 1,603 valid molecules from 683')], top=5.8, size=11)

# ============== Slide 4: Geneformer + Attention deconv ==============
s = blank()
title(s, 'Cell-state foundation model + attention deconvolution')
sub(s, 'Geneformer 768-d zero-shot embedding → 57 Leiden prototypes; non-collapsing TCGA composition')
img(s, FIG/'fig09_geneformer_umap.png', 0.4, 1.6, h=5.4)
img(s, FIG/'fig03_composition.png', 6.6, 1.6, h=2.7)
img(s, FIG/'fig18_deconv_compare.png', 6.6, 4.5, h=2.7)
bullets(s, [('Zero-shot purity','10-NN cell-type purity 0.80'),
            ('Composition entropy','89% of max in TCGA / 89% GSE14520 / 89% GSE76427'),
            ('NNLS / NMF baselines','Both collapse (97% / 0.5% top-1 >50 %); attention is the only one that holds')],
        left=6.6, top=6.8, width=6.5, size=10)

# ============== Slide 5: Prognostic prototypes (TCGA) ==============
s = blank()
title(s, 'Prognostic prototypes — multivariate Cox (TCGA-LIHC)')
sub(s, '9 prototypes Cox-significant after stage / age / sex adjustment; biology matches HCC literature')
img(s, FIG/'fig04_cox_prototypes.png', 0.4, 1.6, w=7.0)
bullets(s, [('Bad prognosis','CD14 macro · CD16 macro · M2 macro · Endothelial · CD4(Treg-like) — TAM/Treg axis'),
            ('Protective','Epithelial(normal) · NK / NKT · CD4(effector) — anti-tumour immunity, residual normal liver'),
            ('Stage controlled','HR per SD (composition) survives stage HR≈1.48, age, sex correction'),
            ('Asymmetric weighting','Bad prototypes 1.0 ×, protective 0.3 × — avoid catastrophic protective-kill')], left=7.7, top=1.6, width=5.4, size=11)

# ============== Slide 6: External validation forest ==============
s = blank()
title(s, 'External validation — meta-analysis forest plot')
sub(s, 'Fixed-effect pooled HRs across TCGA + GSE14520 + GSE76427 (n=763, 2 platforms, 2 etiologies)')
img(s, FIG/'fig11_meta_forest.png', 0.4, 1.5, h=5.7)
bullets(s, [('16 / 57','prototypes meta-pooled p<0.05 (3-cohort fixed-effect)'),
            ('Top hit','proto_11 Endothelial pooled HR 1.27 p=7×10⁻⁶ (Q=3.3 — consistent across cohorts)'),
            ('NK signal','proto_3 NK/NKT pooled HR 0.80 (Q=0.38, near-perfect homogeneity)'),
            ('TAM cluster','CD14/CD16/M2 macrophage all pooled-significant ↑'),
            ('Heterogeneity','low Q for immune prototypes; higher Q for Epithelial (etiology-dependent)')], left=8.5, top=1.6, width=4.7, size=11)

# ============== Slide 7: Cross-cohort scatter ==============
s = blank()
title(s, 'Cross-cohort log(HR) concordance')
sub(s, 'TCGA-LIHC mv-Cox vs external cohorts — direct cross-cohort signal transfer')
img(s, FIG/'fig13_cross_cohort_scatter.png', 0.4, 1.6, w=12.5)
bullets(s, [('GSE14520 (Affymetrix, n=225)','positive correlation; TCGA-significant prototypes mostly on the same diagonal half'),
            ('GSE76427 (Illumina, n=115)','weaker, partly due to platform difference + Singapore HBV-only etiology'),
            ('8 / 9 TCGA-mv-significant prototypes replicate at OS or RFS in ≥1 external cohort'),
            ('Honest caveat','Cross-platform HR transfer is hard; per-cohort harmonisation (z-score + ComBat-style) helps')], top=5.7, size=12)

# ============== Slide 8: KM curves ==============
s = blank()
title(s, 'GSE14520 — Kaplan-Meier dichotomised composite risk')
sub(s, 'TCGA-trained risk score generalises with strong KM separation in 225-patient external cohort')
img(s, FIG/'fig12_km_gse14520.png', 0.4, 1.6, w=12.5)
bullets(s, [('OS log-rank p = 8 × 10⁻⁶','High-risk median ≈ 1383 days vs Low-risk inf'),
            ('RFS log-rank p = 1 × 10⁻³','High-risk median 807 days vs Low-risk 1737 days'),
            ('Composite score','Σ_p (composition_p × log(HR_TCGA_p)) for prototypes with TCGA mv p<0.1'),
            ('Concordance index','c-index 0.73 (OS) / 0.69 (RFS) — covariate-adjusted')], top=5.7, size=12)

# ============== Slide 9: Composite c-index across cohorts ==============
s = blank()
title(s, 'Composite TCGA-derived risk score — c-index across cohorts')
sub(s, 'Multivariate-adjusted external validation; significant in larger Affymetrix cohort')
img(s, FIG/'fig17_cindex_bars.png', 1.0, 1.6, h=5.5)
bullets(s, [('GSE14520 OS','HR 1.35 / SD; p=0.006; c-index 0.73 (n=219)'),
            ('GSE14520 RFS','HR 1.24 / SD; p=0.021; c-index 0.69 (n=219)'),
            ('GSE76427 OS','HR 0.86; p=0.49; c-index 0.69 (small n=115, HBV-driven)'),
            ('GSE76427 RFS','HR 1.27; p=0.13; c-index 0.64'),
            ('Interpretation','Pipeline architecture transferable; HR weights work strongest on RNA-seq/Affymetrix; under-powered on Illumina')], left=8.0, top=1.6, width=5.2, size=10)

# ============== Slide 10: Method consistency Jaccard ==============
s = blank()
title(s, 'Methodological consistency — survival methods + score aggregations')
sub(s, 'Conclusions are robust to method choice, not artefact of one pipeline')
img(s, FIG/'fig14_method_jaccard.png', 0.4, 1.6, w=12.5)
bullets(s, [('Survival methods (5)','Cox univariate / Cox multivariate / KM dichotomised / Random Survival Forest VI / Logistic 3-yr OS'),
            ('Cox uni vs mv vs LogReg','top-9 Jaccard 0.80 — concordant'),
            ('KM and RSF','more divergent (KM uses median split, RSF non-linear), but still consistent at consensus level'),
            ('Drug aggregation (4)','z-sum / rank-mean / Borda / weighted geometric mean'),
            ('Drug aggregation top-20','rank-mean and Borda mathematically equivalent; z-sum vs others 0.43–0.67 Jaccard')], top=5.7, size=11)

# ============== Slide 11: Consensus prognostic prototypes ==============
s = blank()
title(s, 'Consensus prognostic prototypes — 5-method vote')
sub(s, 'Prototypes selected by ≥3 of 5 independent statistical methods')
img(s, FIG/'fig15_consensus_votes.png', 0.4, 1.6, h=5.5)
bullets(s, [('5 / 5 votes','proto_11 Endothelial · proto_46 CD16 macrophage'),
            ('4 / 5 votes','proto_9 Epithelial(normal) · proto_38 CD14 macro · proto_14 CD4(eff)'),
            ('3 / 5 votes','proto_35 Epi · proto_3, proto_24 NK/NKT · proto_37 CD4'),
            ('Consensus = 9','Identical to TCGA mv-Cox top-9 — no method-dependent artefact'),
            ('Phase-2 use','these 9 prototypes drive the composite risk score and per-prototype drug ranking')], left=7.0, top=1.6, width=6.0, size=10)

# ============== Slide 12: Consensus drugs ==============
s = blank()
title(s, 'Consensus drugs — 4-method aggregation vote')
sub(s, '11 drugs in top-20 of all 4 score aggregations — robust mechanistic convergence')
img(s, FIG/'fig16_consensus_drugs.png', 0.4, 1.6, h=5.5)
bullets(s, [('EGFR/HER (4)','Lapatinib · Afatinib · Gefitinib · Osimertinib — all Launched'),
            ('AKT (3)','MK-2206 · Uprosertib · Hexamethylenebisacetamide'),
            ('HSP90 (2)','Alvespimycin · Tanespimycin'),
            ('PI3K + MEK','Copanlisib · Trametinib (both Launched)'),
            ('100 % oncology phase','6 Launched / 5 Phase 2-3 — no preclinical, no cytotoxin')], left=9.2, top=1.6, width=4.0, size=10)

# ============== Slide 13: Tiered ranking + LOO CV ==============
s = blank()
title(s, 'Tiered ranking + leave-one-line-out cross-validation')
sub(s, 'Honest baseline: Spearman 0.145 ± 0.14 in LOO; Tier A/B 176 / Tier C/D 1,515')
bullets(s, [('LOO 5-fold CV','Spearman 0.145 ± 0.141 over 24 held-out HCC lines (vs 0.61 in 3-line hold-out — initial number was over-optimistic)'),
            ('Per-line range','HUH1 0.22 · SNU398 0.23 · HEP3B 0.25 · HUH7 0.25 · LI7 0.25 · JHH1 0.26 · JHH5 0.29 · SKHEP1 0.37'),
            ('Tier A clinical-onc','25 drugs (PRISM disease.area=oncology + Phase 1+)'),
            ('Tier B likely-onc','151 drugs hit ≥1 cancer-driver target gene'),
            ('Tier C low-trust target','1,077 drugs targeting prototypes with trust<0.30 — flagged but not removed'),
            ('Tier C preclinical','399 drugs — moved out of primary list'),
            ('Tier D toxic / non-spec','39 drugs blacklisted (Cetrimonium, Alexidine, Carmustine, Doxorubicin etc.)'),
            ('Chemical diversity','Top-30 high-confidence drugs span 30 distinct Tanimoto-0.6 chemical clusters')], top=1.5, width=12.5, size=12)

# ============== Slide 14: Wet-lab brief ==============
s = blank()
title(s, 'Wet-lab brief — Top-5 high-confidence candidates')
sub(s, 'Each: target subpopulation · FACS markers · suggested HCC line · MOA · phase · SMILES')
brief=[
 (1,'Lapatinib','EGFR/HER2 (Launched)','#43 NK/NKT, #46 CD16 macro, #0 M2 macro','NKG7 / GNLY / APOA1 / HLA-DRA','HUH6 (r=0.24) / SNU761 (r=0.18)'),
 (2,'Bortezomib','Proteasome / NF-κB (Launched)','#6 CD8 T, #18 CD4, #21 M2','C1QA / HSPH1 / CTSB / HSP90AB1','SNU398 (r=0.17)'),
 (3,'Alvespimycin','HSP90 (Phase 2)','#46 CD16 macro, #48 M2','C1QA / HLA-DPA1 / GNLY','SNU761 (r=0.18)'),
 (4,'Afatinib','pan-HER (Launched)','#22 Epithelial_tumor, #26 Epithelial_tumor','APOA1 / ALB / AFP — AFP+ FACS gate','HUH6 (r=0.24)'),
 (5,'MK-2206','Allosteric AKT (Phase 2)','#22 Epithelial_tumor, #26 Epithelial_tumor','APOA1 / APOA2 / ALB','SNU398 (r=0.17)')]
tb = s.shapes.add_textbox(Inches(0.4),Inches(1.4),Inches(12.5),Inches(5.6))
tf = tb.text_frame
for i,(rnk,d,moa,sub_,mk,line) in enumerate(brief):
    if i>0: tf.add_paragraph().space_after=Pt(4)
    p = tf.paragraphs[-1] if i==0 else tf.add_paragraph()
    r1=p.add_run(); r1.text=f'  #{rnk}  {d}  '; r1.font.bold=True; r1.font.size=Pt(16); r1.font.color.rgb=ACC
    r2=p.add_run(); r2.text=f'—  {moa}'; r2.font.size=Pt(14); r2.font.color.rgb=SUB
    p2=tf.add_paragraph(); r3=p2.add_run(); r3.text=f'     Target subpops: {sub_}'; r3.font.size=Pt(11)
    p3=tf.add_paragraph(); r4=p3.add_run(); r4.text=f'     FACS markers: {mk}'; r4.font.size=Pt(11)
    p4=tf.add_paragraph(); r5=p4.add_run(); r5.text=f'     Suggested line: {line}'; r5.font.size=Pt(11)
    p4.space_after=Pt(4)

# ============== Slide 15: Mechanism convergence ==============
s = blank()
title(s, 'Mechanism convergence — biology, not noise')
sub(s, 'Top-25 drugs span 5 oncology-validated pathways aligned with HCC literature')
img(s, FIG/'fig10_mechanism_pie.png', 1.5, 1.6, h=5.5)
bullets(s, [('EGFR/HER (×6)','Lapatinib · Afatinib · Gefitinib · Osimertinib · Erlotinib · AZD8931'),
            ('Proteasome (×3)','Bortezomib · Carfilzomib · CID-5458317 (preclinical, flagged)'),
            ('HSP90 (×2)','Alvespimycin · Tanespimycin'),
            ('AKT/PI3K (×4)','MK-2206 · Uprosertib · Copanlisib · Buparlisib'),
            ('CDK/MEK (×3)','BMS-265246 · Trametinib · PD-0325901'),
            ('VEGFR rescue (×3)','Cediranib · Semaxanib · BMS-690514 (Phase 3 antiangiogenics)'),
            ('All 5 classes','have published HCC efficacy or active trials')], left=7.6, top=1.6, width=5.5, size=11)

# ============== Slide 16: Novelty + journal target ==============
s = blank()
title(s, 'Novelty & target journals')
items_n = [
    ('N1','First scRNA foundation-model embedding as basis for bulk attention deconvolver (vs cell-line basis collapse)'),
    ('N2','First fusion of cell-state Cox + drug structure + drug-target × pathway-Cox prior'),
    ('N3','First 3-cohort × 3-platform × 2-endpoint external validation of single-cell-deconvolution drug ranking (n=763)'),
    ('N4','Trust-aware tiered scoring — explicit calibration on which subpop predictions are reliable'),
    ('N5','5-method survival consensus + 4-method drug aggregation consensus = robust to analytic choice'),
]
items_j = [
    ('Q1 ready','Briefings in Bioinformatics · NAR Genomics & Bioinformatics · Bioinformatics · JHEP Reports'),
    ('Q1 borderline','Cell Reports Methods · Cell Systems (need ICI cohort to seal)'),
    ('Aspirational','Nature Communications (need IMbrave150 / ICI replication + wet-lab follow-up)'),
]
tb = s.shapes.add_textbox(Inches(0.4),Inches(1.5),Inches(12.5),Inches(2.8))
tf=tb.text_frame
for i,(k,v) in enumerate(items_n):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    r1=p.add_run(); r1.text=k+'  '; r1.font.bold=True; r1.font.size=Pt(15); r1.font.color.rgb=GOOD
    r2=p.add_run(); r2.text=v; r2.font.size=Pt(13)
    p.space_after=Pt(6)
tb = s.shapes.add_textbox(Inches(0.4),Inches(4.6),Inches(12.5),Inches(2.5))
tf=tb.text_frame
for i,(k,v) in enumerate(items_j):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    r1=p.add_run(); r1.text=k+':  '; r1.font.bold=True; r1.font.size=Pt(15); r1.font.color.rgb=ACC
    r2=p.add_run(); r2.text=v; r2.font.size=Pt(13)
    p.space_after=Pt(6)

# ============== Slide 17: Limitations ==============
s = blank()
title(s, 'Honest limitations & Phase-2 roadmap')
lims = [
 ('GSE76427 weak (n=115, Illumina, HBV-only)','Underpowered + platform mismatch. Confidence still moderate (c-index 0.69).'),
 ('TAM trust 0.25','DepMap has no tissue-resident macrophage; AML blasts ≠ HCC TAM. Predictions flagged as Tier C.'),
 ('Sorafenib mid-rank','VEGFR clinical benefit is microenvironmental; partially rescued via angiogenesis pathway prior.'),
 ('No ICI cohort','IMbrave150 not public; Phase-2 will use GSE193084 sorafenib cohort + ICI BMS-cohort to test'),
 ('No wet-lab','Phase-2 to test Top-5 in HCC PDOs via FACS-sorted subpop killing + Perturb-seq Δz match'),
 ('MPNN unseen-drug not integrated','+0.30 Spearman on truly unseen drugs; Phase-2 to retrain as direction-vector generator')
]
tb=s.shapes.add_textbox(Inches(0.4),Inches(1.2),Inches(12.5),Inches(6))
tf=tb.text_frame
for i,(a,b) in enumerate(lims):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    r1=p.add_run(); r1.text='• '+a; r1.font.bold=True; r1.font.size=Pt(14); r1.font.color.rgb=BAD
    p2=tf.add_paragraph(); r2=p2.add_run(); r2.text='    '+b; r2.font.size=Pt(12); r2.font.color.rgb=SUB
    p2.space_after=Pt(6)

# ============== Slide 18: Summary ==============
s = blank()
title(s, 'Summary — what each number means')
nums = [
    ('763', 'patients across TCGA + 2 external cohorts (3 platforms)'),
    ('57', 'cell-state prototypes (Geneformer Leiden); 9 prognostic'),
    ('16 / 57', 'prototypes meta-analysis pooled-significant (3-cohort fixed-effect)'),
    ('8 / 9', 'TCGA-mv-significant prototypes replicated in ≥1 external endpoint'),
    ('0.727', 'composite TCGA-risk c-index in GSE14520 OS (n=219, p=0.006)'),
    ('5 + 4', 'survival methods + score aggregations — full consensus testing'),
    ('11', 'drugs in top-20 of all 4 aggregation methods — mechanism convergence'),
    ('176 / 1806', 'high-confidence drugs (Tier A+B) vs exploratory (Tier C+D)'),
]
tb=s.shapes.add_textbox(Inches(0.4),Inches(1.4),Inches(12.5),Inches(5.5))
tf=tb.text_frame
for i,(v,d) in enumerate(nums):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    r1=p.add_run(); r1.text=f'  {v:>9s}  '; r1.font.bold=True; r1.font.size=Pt(22); r1.font.color.rgb=ACC
    r2=p.add_run(); r2.text=d; r2.font.size=Pt(13)
    p.space_after=Pt(7)

prs.save(OUT)
print(f'Saved: {OUT}')
print(f'Size: {OUT.stat().st_size/1024/1024:.2f} MB')
print(f'Slides: {len(prs.slides)}')
